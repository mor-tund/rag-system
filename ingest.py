"""
Pipeline ingestion cho RAG demo (Bước 4).
Đọc proposal (xlsx) + case study (pptx) trong data/ → nạp vào Postgres:
  - opportunity / opportunity_wbs_item : dữ liệu CÓ CẤU TRÚC (SQL)
  - case_study                         : metadata dự án đã làm
  - document_chunk                     : text chunk + vector (bge-m3, local)

Chạy:  .venv/bin/python ingest.py
"""
import os
import re
import glob

import psycopg
from pgvector.psycopg import register_vector
from openpyxl import load_workbook
from pptx import Presentation
from sentence_transformers import SentenceTransformer

# ----------------------------- Cấu hình -----------------------------
DB_DSN = os.environ.get("RAG_DB_DSN", "postgresql://rag:ragpass@localhost:5433/rag")
EMBED_MODEL = "BAAI/bge-m3"          # đa ngữ (VI/KO/JA/EN), vector 1024 chiều, chạy local
EMBED_DIM = 1024
HERE = os.path.dirname(os.path.abspath(__file__))
DATA_DIR = os.path.join(HERE, "data")
CHUNK_MAX = 1600                      # giới hạn ký tự mỗi chunk

# ----------------------------- Embedding -----------------------------
_model = None
def get_model():
    global _model
    if _model is None:
        print(f"[embed] Đang nạp model {EMBED_MODEL} (lần đầu sẽ tải ~2.3GB)...", flush=True)
        _model = SentenceTransformer(EMBED_MODEL)
        print("[embed] Model sẵn sàng.", flush=True)
    return _model

def embed(texts):
    model = get_model()
    return model.encode(texts, normalize_embeddings=True, show_progress_bar=False)

# ----------------------------- Tiện ích parse -----------------------------
def _norm(v):
    return str(v).strip() if v is not None else ""

def _num(v):
    """Ép về float nếu được, ngược lại None."""
    if v is None:
        return None
    try:
        return float(str(v).replace(",", "").strip())
    except (ValueError, TypeError):
        return None

def _find_after(rows, label, max_gap=5):
    """Tìm ô bắt đầu bằng `label`, trả giá trị ô không rỗng kế tiếp cùng dòng."""
    for r in rows:
        for ci, v in enumerate(r):
            if _norm(v).startswith(label):
                for nv in r[ci + 1: ci + 1 + max_gap]:
                    if _norm(nv):
                        return _norm(nv)
    return None

def chunk_text(text, max_len=CHUNK_MAX):
    """Cắt text dài thành các đoạn <= max_len, ưu tiên ranh giới dòng trống."""
    text = text.strip()
    if len(text) <= max_len:
        return [text] if text else []
    chunks, buf = [], ""
    for para in re.split(r"\n\s*\n", text):
        if len(buf) + len(para) + 2 <= max_len:
            buf = (buf + "\n\n" + para).strip()
        else:
            if buf:
                chunks.append(buf)
            if len(para) <= max_len:
                buf = para
            else:  # đoạn quá dài: cắt cứng
                for i in range(0, len(para), max_len):
                    chunks.append(para[i:i + max_len])
                buf = ""
    if buf:
        chunks.append(buf)
    return chunks

# ----------------------------- Parse proposal (xlsx) -----------------------------
NUM_RE = re.compile(r"^\d+(\.\d+)?$")

def parse_proposal(path):
    wb = load_workbook(path, data_only=True)
    overall = list(wb["Overall"].iter_rows(values_only=True)) if "Overall" in wb.sheetnames else []
    assum = list(wb["Assumption"].iter_rows(values_only=True)) if "Assumption" in wb.sheetnames else []
    wbs_sheet = next((s for s in wb.sheetnames if s.startswith("WBS Function List (MOR)")), None)
    wbs_rows = list(wb[wbs_sheet].iter_rows(values_only=True)) if wbs_sheet else []

    tech = _find_after(assum, "Frontend") or _find_after(assum, "Language/ Development")
    if tech:
        tech = tech.replace("\n", ", ").strip()

    header = {
        "name": _find_after(overall, "Project name") or "HCMS",
        "customer": _find_after(overall, "Customer") or "HiveLab",
        "department": _find_after(overall, "Department") or "MSOL",
        "doc_type": "estimate",
        "tech_stack": tech or "Frontend: React JS, Backend: Java",
        "total_effort_mm": _num(_find_after(overall, "Total")) or 18.034,
        "total_effort_md": _num(_find_after(wbs_rows, "Total（MD）")) or 239.52,
        "timeline_months": 4,
        "language": _find_after(assum, "Language Support") or "English",
        "source_date": "2026-06-04",
        "status": "draft",
        "owner": None,
        "description": "Đề xuất/ước lượng dự án HCMS (hệ thống quản lý hợp đồng) cho khách hàng HiveLab. "
                       "Frontend React JS, Backend Java. Gồm cải tiến quản lý hợp đồng, quyết toán, đăng nhập, bảo mật.",
    }

    # WBS items
    items, category = [], None
    for r in wbs_rows:
        c1 = _norm(r[1]) if len(r) > 1 else ""
        if not NUM_RE.match(c1):
            continue
        cat = _norm(r[2]) if len(r) > 2 else ""
        if cat and not NUM_RE.match(cat):
            category = cat
        name = (_norm(r[3]) if len(r) > 3 else "") or (_norm(r[4]) if len(r) > 4 else "") or category
        desc = _norm(r[5]) if len(r) > 5 else ""
        note = _norm(r[12]) if len(r) > 12 else ""
        m = re.search(r"Độ quan trọng[:\s]*([^\n【]+)", note)
        priority = m.group(1).strip() if m else None
        items.append({
            "category": category, "name": name, "description": desc,
            "effort_study": _num(r[6]) if len(r) > 6 else None,
            "effort_fe": _num(r[7]) if len(r) > 7 else None,
            "effort_be": _num(r[8]) if len(r) > 8 else None,
            "effort_ut": _num(r[9]) if len(r) > 9 else None,
            "effort_total": _num(r[10]) if len(r) > 10 else None,
            "priority": priority,
        })

    # Các khối Assumption/scope (mỗi dòng có nội dung = 1 khối để embed)
    assum_blocks = []
    for r in assum:
        cells = [_norm(c) for c in r if _norm(c)]
        if not cells:
            continue
        text = " | ".join(cells)
        if len(text) >= 25 and not text.startswith("Assumption | Project name"):
            assum_blocks.append(text)

    # Khối Overall: phân bổ effort theo loại việc + resource plan theo tháng
    overall_lines = []
    for r in overall:
        cells = [_norm(c) for c in r if _norm(c)]
        if cells:
            overall_lines.append(" | ".join(cells))
    overall_text = ("Overall — Phân bổ effort theo loại việc và Kế hoạch nhân sự "
                    "(Resource Plan) theo tháng của dự án HCMS:\n" + "\n".join(overall_lines))
    overall_blocks = chunk_text(overall_text)

    return header, items, assum_blocks, overall_blocks

# ----------------------------- Parse case study (pptx) -----------------------------
TECH_KEYWORDS = [
    "React Native", "React", "Node.js", "Express", "Java", "Spring Boot", "MySQL",
    "MongoDB", "Redis", "Stripe", "Socket.io", "Azure", "AWS", "Expo", "BullMQ",
    "Sequelize", "Firebase", "Twilio", "SendGrid", "Vite", "Flutter", "Kotlin",
    "Swift", "PostgreSQL", "Next.js", "Django", "Python", "GraphQL", "Docker",
]
DOMAIN_MAP = [
    ("driver", "Ride-hailing / Đặt xe"), ("ride", "Ride-hailing / Đặt xe"),
    ("insurance", "Bảo hiểm"), ("logistics", "Logistics / WMS"), ("wms", "Logistics / WMS"),
    ("epump", "Bán lẻ xăng dầu / ePump"), ("music", "Sự kiện âm nhạc"),
    ("conference", "Hội nghị / Conference"), ("glueup", "Quản lý sự kiện / cộng đồng"),
]

def parse_pptx(path):
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    cells = " | ".join(c.text.strip() for c in row.cells)
                    if cells.strip(" |"):
                        parts.append(cells)
        txt = "\n".join(parts).strip()
        if txt:
            slides.append((i, txt))
    return slides

def case_study_meta(path, slides):
    fname = os.path.basename(path).lower()
    full = "\n".join(t for _, t in slides)
    title = slides[0][1].split("\n")[1] if slides and len(slides[0][1].split("\n")) > 1 else os.path.basename(path)
    techs = [k for k in TECH_KEYWORDS if re.search(r"\b" + re.escape(k) + r"\b", full, re.I)]
    domain = next((d for key, d in DOMAIN_MAP if key in fname), None)
    customer = "Ẩn theo NDA"
    for brand in ["FWD", "GlueUp", "Glue Up"]:
        if re.search(re.escape(brand), full, re.I):
            customer = brand
            break
    return {
        "name": os.path.splitext(os.path.basename(path))[0],
        "title": title.strip()[:300],
        "customer": customer,
        "domain": domain,
        "tech_stack": ", ".join(dict.fromkeys(techs)) or None,
        "file_path": os.path.relpath(path, HERE),
    }

# ----------------------------- Nạp vào DB -----------------------------
def main():
    conn = psycopg.connect(DB_DSN)
    register_vector(conn)
    cur = conn.cursor()

    print("[db] Xoá dữ liệu cũ (nếu có) để nạp lại sạch...", flush=True)
    cur.execute("TRUNCATE document_chunk, opportunity_wbs_item, opportunity, case_study RESTART IDENTITY CASCADE;")

    pending = []  # (source_type, source_id, chunk_index, content, metadata)

    # ---- Proposal → opportunity ----
    prop_files = sorted(glob.glob(os.path.join(DATA_DIR, "proposal", "*.xlsx")))
    for path in prop_files:
        print(f"[opp] Parse proposal: {os.path.basename(path)}", flush=True)
        header, items, assum_blocks, overall_blocks = parse_proposal(path)
        cur.execute(
            """INSERT INTO opportunity
               (name, customer, department, doc_type, tech_stack, total_effort_mm,
                total_effort_md, timeline_months, language, source_date, status, owner, description)
               VALUES (%(name)s,%(customer)s,%(department)s,%(doc_type)s,%(tech_stack)s,
                       %(total_effort_mm)s,%(total_effort_md)s,%(timeline_months)s,%(language)s,
                       %(source_date)s,%(status)s,%(owner)s,%(description)s)
               RETURNING id""", header)
        opp_id = cur.fetchone()[0]
        print(f"      opportunity id={opp_id}  '{header['name']}'  ({len(items)} chức năng WBS)", flush=True)

        for it in items:
            it2 = dict(it, opportunity_id=opp_id)
            cur.execute(
                """INSERT INTO opportunity_wbs_item
                   (opportunity_id, category, name, description, effort_study, effort_fe,
                    effort_be, effort_ut, effort_total, priority)
                   VALUES (%(opportunity_id)s,%(category)s,%(name)s,%(description)s,%(effort_study)s,
                           %(effort_fe)s,%(effort_be)s,%(effort_ut)s,%(effort_total)s,%(priority)s)""", it2)

        # chunks cho opportunity: 1 chunk tổng quan + mỗi chức năng WBS 1 chunk
        idx = 0
        summary = (f"Dự án {header['name']} cho khách hàng {header['customer']} (bộ phận {header['department']}). "
                   f"Loại: {header['doc_type']}. Tech: {header['tech_stack']}. "
                   f"Tổng effort {header['total_effort_mm']} MM / {header['total_effort_md']} MD, "
                   f"thời gian {header['timeline_months']} tháng. {header['description']}")
        pending.append(("opportunity", opp_id, idx, summary,
                        {"file": os.path.basename(path), "kind": "summary"})); idx += 1
        for it in items:
            block = f"[{it['category']}] {it['name']}\n{it['description']}\nEffort total: {it['effort_total']} MD, ưu tiên: {it['priority']}"
            for c in chunk_text(block):
                pending.append(("opportunity", opp_id, idx, c,
                                {"file": os.path.basename(path), "kind": "wbs_item",
                                 "category": it["category"], "name": it["name"]})); idx += 1
        # chunks cho phần Assumption/scope của opportunity
        for blk in assum_blocks:
            for c in chunk_text(blk):
                pending.append(("opportunity", opp_id, idx, c,
                                {"file": os.path.basename(path), "kind": "assumption"})); idx += 1
        print(f"      + {len(assum_blocks)} khối assumption/scope", flush=True)
        # chunks cho phần Overall (effort breakdown + resource plan)
        for blk in overall_blocks:
            pending.append(("opportunity", opp_id, idx, blk,
                            {"file": os.path.basename(path), "kind": "overall"})); idx += 1
        print(f"      + {len(overall_blocks)} khối overall (effort + resource plan)", flush=True)

    # ---- Case studies → case_study ----
    cs_files = sorted(glob.glob(os.path.join(DATA_DIR, "caseStudy", "*.pptx")))
    for path in cs_files:
        slides = parse_pptx(path)
        meta = case_study_meta(path, slides)
        cur.execute(
            """INSERT INTO case_study (name, title, customer, domain, tech_stack, file_path)
               VALUES (%(name)s,%(title)s,%(customer)s,%(domain)s,%(tech_stack)s,%(file_path)s)
               RETURNING id""", meta)
        cs_id = cur.fetchone()[0]
        print(f"[cs ] {meta['name']}  id={cs_id}  domain={meta['domain']}  tech={meta['tech_stack']}", flush=True)
        idx = 0
        for slide_no, txt in slides:
            for c in chunk_text(txt):
                pending.append(("case_study", cs_id, idx, c,
                                {"file": os.path.basename(path), "slide": slide_no})); idx += 1

    # ---- Embed tất cả chunk một lần rồi nạp ----
    print(f"\n[embed] Tổng {len(pending)} chunk cần vector hoá...", flush=True)
    texts = [p[3] for p in pending]
    vectors = embed(texts)
    print("[db] Đang ghi document_chunk...", flush=True)
    import json
    for (stype, sid, cidx, content, meta), vec in zip(pending, vectors):
        cur.execute(
            """INSERT INTO document_chunk (source_type, source_id, chunk_index, content, embedding, metadata)
               VALUES (%s,%s,%s,%s,%s,%s)""",
            (stype, sid, cidx, content, vec, json.dumps(meta, ensure_ascii=False)))

    conn.commit()

    # ---- Báo cáo ----
    print("\n===== KẾT QUẢ NẠP DỮ LIỆU =====", flush=True)
    for tbl in ["opportunity", "opportunity_wbs_item", "case_study", "document_chunk"]:
        cur.execute(f"SELECT count(*) FROM {tbl}")
        print(f"  {tbl:24s}: {cur.fetchone()[0]} dòng", flush=True)
    cur.execute("SELECT source_type, count(*) FROM document_chunk GROUP BY source_type")
    for st, n in cur.fetchall():
        print(f"     - chunk {st}: {n}", flush=True)
    cur.close()
    conn.close()
    print("HOÀN TẤT.", flush=True)

if __name__ == "__main__":
    main()