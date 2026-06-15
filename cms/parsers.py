"""Trích text từ nhiều định dạng tài liệu, trả về danh sách (nhãn, text) sẵn sàng để chunk."""
import os
import re

CHUNK_MAX = 1600
IMG_EXT = {".png", ".jpg", ".jpeg", ".webp", ".tiff", ".bmp"}
SUPPORTED = {".pptx", ".xlsx", ".xls", ".pdf", ".docx", ".txt", ".md"} | IMG_EXT

OCR_LANG = os.environ.get("OCR_LANG", "vie+eng+jpn")  # gói ngôn ngữ tesseract
OCR_MIN_TEXT = 20   # trang có < ngần này ký tự coi như scan -> OCR
OCR_DPI = 200


def _ocr_available() -> bool:
    try:
        import pytesseract  # noqa: F401
        from PIL import Image  # noqa: F401
        return True
    except Exception:
        return False


def chunk_text(text, max_len=CHUNK_MAX):
    text = (text or "").strip()
    if len(text) <= max_len:
        return [text] if text else []
    chunks, buf = [], ""
    for para in re.split(r"\n\s*\n", text):
        if len(buf) + len(para) + 2 <= max_len:
            buf = (buf + "\n\n" + para).strip()
        else:
            if buf:
                chunks.append(buf)
            if len(para) <= max_len:
                buf = para
            else:
                for i in range(0, len(para), max_len):
                    chunks.append(para[i:i + max_len])
                buf = ""
    if buf:
        chunks.append(buf)
    return chunks


def _pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    out = []
    for i, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    cells = " | ".join(c.text.strip() for c in row.cells)
                    if cells.strip(" |"):
                        parts.append(cells)
        txt = "\n".join(parts).strip()
        if txt:
            out.append((f"slide {i}", txt))
    return out


def _xlsx(path):
    from openpyxl import load_workbook
    wb = load_workbook(path, data_only=True)
    out = []
    for ws in wb.worksheets:
        lines = []
        for r in ws.iter_rows(values_only=True):
            cells = [str(c).strip() for c in r if c is not None and str(c).strip()]
            if cells:
                lines.append(" | ".join(cells))
        if lines:
            out.append((f"sheet {ws.title}", "\n".join(lines)))
    return out


def _ocr_pixmap(pix):
    """OCR một pixmap PyMuPDF → text (tesseract, đa ngữ)."""
    import io
    import pytesseract
    from PIL import Image
    img = Image.open(io.BytesIO(pix.tobytes("png")))
    return pytesseract.image_to_string(img, lang=OCR_LANG).strip()


def _pdf(path):
    """PDF: ưu tiên text layer; trang nào thiếu text (scan) → OCR (nếu có tesseract)."""
    try:
        import fitz  # PyMuPDF: vừa trích text vừa render ảnh để OCR
    except Exception:
        # fallback: pypdf text-only, không OCR
        from pypdf import PdfReader
        out = []
        for i, page in enumerate(PdfReader(path).pages, 1):
            txt = (page.extract_text() or "").strip()
            if txt:
                out.append((f"page {i}", txt))
        return out

    out = []
    use_ocr = _ocr_available()
    doc = fitz.open(path)
    for i, page in enumerate(doc, 1):
        txt = (page.get_text() or "").strip()
        if len(txt) < OCR_MIN_TEXT and use_ocr:  # trang scan → OCR
            try:
                txt = _ocr_pixmap(page.get_pixmap(dpi=OCR_DPI)) or txt
            except Exception:
                pass
        if txt:
            out.append((f"page {i}", txt))
    return out


def _image(path):
    """File ảnh → OCR (cần tesseract)."""
    if not _ocr_available():
        raise ValueError("Ảnh cần OCR nhưng chưa cài tesseract/pytesseract trên server.")
    import pytesseract
    from PIL import Image
    txt = pytesseract.image_to_string(Image.open(path), lang=OCR_LANG).strip()
    return [("image", txt)] if txt else []


def _docx(path):
    import docx
    d = docx.Document(path)
    paras = [p.text.strip() for p in d.paragraphs if p.text.strip()]
    for table in d.tables:
        for row in table.rows:
            cells = " | ".join(c.text.strip() for c in row.cells)
            if cells.strip(" |"):
                paras.append(cells)
    text = "\n".join(paras).strip()
    return [("document", text)] if text else []


def _txt(path):
    with open(path, encoding="utf-8", errors="ignore") as f:
        text = f.read().strip()
    return [("document", text)] if text else []


def extract_blocks(path):
    """Trả về list (nhãn vị trí, text). Raise ValueError nếu định dạng chưa hỗ trợ."""
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pptx":
        return _pptx(path)
    if ext in (".xlsx", ".xls"):
        return _xlsx(path)
    if ext == ".pdf":
        return _pdf(path)
    if ext == ".docx":
        return _docx(path)
    if ext in (".txt", ".md"):
        return _txt(path)
    if ext in IMG_EXT:
        return _image(path)
    raise ValueError(f"Định dạng chưa hỗ trợ: {ext} (hỗ trợ: {', '.join(sorted(SUPPORTED))}).")